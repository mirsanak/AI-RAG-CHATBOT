import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import pandas as pd
from pptx import Presentation
from docx import Document
from bs4 import BeautifulSoup
import time
from google.api_core.exceptions import ResourceExhausted

# Hardcoded API Key
api_key = "AIzaSyCovFBjdOAdflN76ozliR09pBclmfeOsE8"

# Set page configuration
st.set_page_config(page_title="Document Genie", layout="wide", page_icon="🤖")

# Custom CSS for colorful and interactive interface
st.markdown("""
<style>
    .stButton button {
        background-color: #4CAF50;
        color: white;
        border-radius: 5px;
        padding: 10px 20px;
        font-size: 16px;
        transition: background-color 0.3s ease;
    }
    .stButton button:hover {
        background-color: #45a049;
    }
    .stTextInput input {
        border-radius: 5px;
        padding: 10px;
        font-size: 16px;
    }
    .stMarkdown h1 {
        color: #4CAF50;
        text-align: center;
    }
    .stMarkdown h2 {
        color: #2E86C1;
    }
    .stMarkdown h3 {
        color: #D35400;
    }
    .stSidebar {
        background-color: #F4F6F6;
        padding: 20px;
        border-radius: 10px;
    }
    .stSpinner div {
        color: #4CAF50;
    }
    .stSuccess {
        color: #28B463;
    }
    .stError {
        color: #E74C3C;
    }
</style>
""", unsafe_allow_html=True)

# App title and description
st.markdown("""
# 🤖 Document Genie: Get Instant Insights from Your Documents

Welcome to **Document Genie**, your AI-powered assistant for extracting insights from documents! This chatbot uses the **Retrieval-Augmented Generation (RAG)** framework, powered by Google's **Gemini-PRO** model, to provide accurate and contextually relevant answers.

### How It Works

1. **Upload Your Documents**: Upload multiple files (PDF, PPTX, Excel, Word, HTML, CSV) for analysis.
2. **Ask a Question**: After processing, ask any question related to the content of your documents.
3. **Get Instant Answers**: Receive precise and detailed answers instantly!

---

""")

# Function to extract text from PDF
def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

# Function to extract text from PPTX
def get_pptx_text(pptx_docs):
    text = ""
    for pptx in pptx_docs:
        prs = Presentation(pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

# Function to extract text from Excel
def get_excel_text(excel_docs):
    text = ""
    for excel in excel_docs:
        df = pd.read_excel(excel)
        text += df.to_string()
    return text

# Function to extract text from Word
def get_word_text(word_docs):
    text = ""
    for word in word_docs:
        doc = Document(word)
        for para in doc.paragraphs:
            text += para.text + "\n"
    return text

# Function to extract text from HTML
def get_html_text(html_docs):
    text = ""
    for html in html_docs:
        soup = BeautifulSoup(html, "html.parser")
        text += soup.get_text()
    return text

# Function to extract text from CSV
def get_csv_text(csv_docs):
    text = ""
    for csv in csv_docs:
        df = pd.read_csv(csv)
        text += df.to_string()
    return text

# Function to split text into chunks
def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

# Function to create a vector store
def get_vector_store(text_chunks, api_key):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=api_key)
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

# Function to create a conversational chain
def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}?\n
    Question: \n{question}\n

    Answer:
    """
    try:
        model = ChatGoogleGenerativeAI(model="gemini-1.5-pro", temperature=0.3, google_api_key=api_key)
        prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
        chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
        return chain
    except ResourceExhausted:
        st.error("API quota exhausted. Please try again later or increase your quota.")
        return None

# Function to handle user input and generate responses
def user_input(user_question, api_key):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", google_api_key=api_key)
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    if chain:
        try:
            response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
            st.markdown(f"**Reply:** {response['output_text']}", unsafe_allow_html=True)
        except ResourceExhausted:
            st.error("API quota exhausted. Please try again later or increase your quota.")
    time.sleep(1)  # Add a delay to avoid rate limiting

# Main function
def main():
    st.markdown("### 📄 Upload Your Documents")
    uploaded_files = st.file_uploader("Choose files (PDF, PPTX, Excel, Word, HTML, CSV)", accept_multiple_files=True, key="file_uploader")

    if st.button("Submit & Process", key="process_button") and api_key:
        with st.spinner("Processing your documents..."):
            raw_text = ""
            for uploaded_file in uploaded_files:
                file_extension = uploaded_file.name.split(".")[-1].lower()
                if file_extension == "pdf":
                    raw_text += get_pdf_text([uploaded_file])
                elif file_extension == "pptx":
                    raw_text += get_pptx_text([uploaded_file])
                elif file_extension in ["xlsx", "xls"]:
                    raw_text += get_excel_text([uploaded_file])
                elif file_extension == "docx":
                    raw_text += get_word_text([uploaded_file])
                elif file_extension == "html":
                    raw_text += get_html_text([uploaded_file])
                elif file_extension == "csv":
                    raw_text += get_csv_text([uploaded_file])
            text_chunks = get_text_chunks(raw_text)
            get_vector_store(text_chunks, api_key)
            st.success("Document processing complete! 🎉")

    st.markdown("### ❓ Ask a Question")
    user_question = st.text_input("Enter your question here:", key="user_question")

    if user_question and api_key:
        user_input(user_question, api_key)

# Run the app
if __name__ == "__main__":
    main()